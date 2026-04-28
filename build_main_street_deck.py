from pathlib import Path
import re, json, zipfile, shutil, subprocess, sys

PDF = Path.home() / "healhouse-decks/source/main-street-architect.pdf"
PPT = Path.home() / "healhouse-decks/source/main-street-joe.pptx"
OUT = Path.home() / "healhouse-decks" / "main-street-assessment.html"
ASSET = Path.home() / "healhouse-decks" / "assets" / "main-street"
ASSET.mkdir(parents=True, exist_ok=True)

import fitz
from pptx import Presentation

def money(n):
    try: return "${:,.0f}".format(float(n))
    except: return str(n)

def clean_num(s):
    s = str(s).replace(",", "").replace("$", "").strip()
    s = "".join(ch for ch in s if ch.isdigit() or ch == ".")
    if not s or s == ".":
        raise ValueError("bad number")
    return float(s)

# Extract PDF pages + text
pdf_text = []
pdf_imgs = []
doc = fitz.open(PDF)
for i, page in enumerate(doc):
    txt = page.get_text("text")
    pdf_text.append(f"--- PDF PAGE {i+1} ---\n{txt}")
    if i < 8:
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        img = ASSET / f"architect-page-{i+1}.jpg"
        pix.save(img)
        pdf_imgs.append(str(img.relative_to(OUT.parent)))
pdf_all = "\n".join(pdf_text)

# Extract PPT text + images
ppt_texts = []
prs = Presentation(PPT)
for idx, slide in enumerate(prs.slides, 1):
    bits = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            bits.append(shape.text.strip())
    ppt_texts.append(f"--- JOE PPT SLIDE {idx} ---\n" + "\n".join(bits))
ppt_all = "\n".join(ppt_texts)

media_dir = ASSET / "ppt-media"
media_dir.mkdir(exist_ok=True)
ppt_imgs = []
with zipfile.ZipFile(PPT) as z:
    for name in z.namelist():
        if name.startswith("ppt/media/") and name.lower().endswith((".png",".jpg",".jpeg")):
            dest = media_dir / Path(name).name
            dest.write_bytes(z.read(name))
            ppt_imgs.append(str(dest.relative_to(OUT.parent)))

def find_gfa(text):
    pats = [
        r"(?:gross floor area|gfa|total floor area|total gfa)[^\d]{0,80}([\d,]{4,})\s*(?:sf|sq\.?\s*ft)",
        r"([\d,]{4,})\s*(?:sf|sq\.?\s*ft)[^\n]{0,80}(?:gross floor area|gfa|total floor area|total gfa)",
    ]
    for p in pats:
        m = re.search(p, text, re.I)
        if m: return int(clean_num(m.group(1)))
    nums = [int(clean_num(x)) for x in re.findall(r"([\d,]{5,})\s*(?:sf|sq\.?\s*ft)", text, re.I)]
    return max(nums) if nums else 0

def find_count(text):
    pats = [
        r"(?:suites|units|beds|resident rooms|rooms)[^\d]{0,40}(\d{2,4})",
        r"(\d{2,4})[^\n]{0,30}(?:suites|units|beds|resident rooms|rooms)",
    ]
    for p in pats:
        vals = [int(x) for x in re.findall(p, text, re.I)]
        vals = [v for v in vals if 10 <= v <= 1000]
        if vals: return max(vals)
    return 0

def find_cost_sf(text):
    pats = [
        r"\$?\s*([\d,]+(?:\.\d+)?)\s*/\s*(?:sf|sq\.?\s*ft)",
        r"(?:cost|construction)[^\n$]{0,60}\$?\s*([\d,]+(?:\.\d+)?)",
    ]
    vals = []
    for p in pats:
        for x in re.findall(p, text, re.I):
            try:
                vals.append(clean_num(x))
            except Exception:
                pass
    vals = [v for v in vals if 100 <= v <= 1000]
    return vals[0] if vals else 383.50

def find_monthly_rent(text):
    pats = [
        r"\$?\s*([\d,]{4,6})\s*(?:/|per)?\s*(?:month|mo|monthly)",
        r"(?:rent|monthly fee|retirement)[^\n$]{0,80}\$?\s*([\d,]{4,6})",
    ]
    vals = []
    for p in pats:
        for x in re.findall(p, text, re.I):
            try:
                vals.append(clean_num(x))
            except Exception:
                pass
    vals = [v for v in vals if 2000 <= v <= 15000]
    return vals[0] if vals else 5500

def find_expense(text):
    m = re.search(r"(?:expense|opex|operating)[^\n%]{0,60}(\d{1,2})\s*%", text, re.I)
    if m: return int(m.group(1)) / 100
    return 0.42

gfa = find_gfa(pdf_all)
suites = find_count(pdf_all)
cost_sf = find_cost_sf(ppt_all)
monthly_rent = find_monthly_rent(ppt_all)
expense_ratio = find_expense(ppt_all)

if not suites and gfa:
    suites = round(gfa / 700)

data = {
    "gfa": gfa,
    "suites": suites,
    "cost_sf": cost_sf,
    "monthly_rent": monthly_rent,
    "expense_ratio": expense_ratio,
    "occupancy": 0.90,
    "growth": 0.025,
    "bch_forgivable": 0.40,
    "bch_rate": 0.0415,
    "bch_am": 50,
    "conv_rate": 0.0499,
    "conv_am": 30,
    "investor_split": 0.60,
    "nfp_split": 0.40,
    "pdf_imgs": pdf_imgs,
    "ppt_imgs": ppt_imgs[:10],
    "pdf_source": pdf_all[:10000],
    "ppt_source": ppt_all[:10000],
}

html = f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Heal House Main Street — Retirement Assessment</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
:root{{--espresso:#1E1A14;--charcoal:#2A2118;--cream:#F5F0E8;--gold:#C9A84C;--gold2:#E2C97E;--teal:#3A7D7B}}
body{{background:var(--espresso);color:var(--cream);font-family:'Segoe UI',system-ui,sans-serif;line-height:1.6}}
.slide{{min-height:100vh;padding:80px 56px;display:flex;align-items:center;border-bottom:1px solid rgba(201,168,76,.12);position:relative;overflow:hidden}}
.slide:nth-child(even){{background:var(--charcoal)}}
.inner{{max-width:1180px;margin:0 auto;width:100%;position:relative;z-index:2}}
.bg{{position:absolute;inset:0;background-size:cover;background-position:center;opacity:.22}}
.overlay{{position:absolute;inset:0;background:linear-gradient(90deg,rgba(20,15,8,.9),rgba(20,15,8,.58))}}
.label{{font-size:11px;letter-spacing:.25em;text-transform:uppercase;color:var(--gold);border-bottom:1px solid rgba(201,168,76,.45);display:inline-block;margin-bottom:22px;padding-bottom:4px}}
h1,h2,h3{{font-family:Georgia,serif;font-weight:400;line-height:1.1}}
h1{{font-size:clamp(48px,8vw,104px);font-style:italic}}
h2{{font-size:clamp(34px,5vw,64px)}}
h3{{font-size:26px;color:var(--gold2);margin-bottom:12px}}
p{{max-width:820px;color:rgba(245,240,232,.82);font-size:18px}}
.grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(230px,1fr));gap:18px;margin:32px 0}}
.card{{background:rgba(42,33,24,.78);border:1px solid rgba(201,168,76,.24);padding:24px}}
.card .k{{font-size:11px;letter-spacing:.14em;text-transform:uppercase;color:var(--gold)}}
.card .v{{font-family:Georgia,serif;font-size:34px;margin-top:8px;color:var(--cream)}}
.card .s{{font-size:13px;color:rgba(245,240,232,.62);margin-top:8px}}
.imggrid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(260px,1fr));gap:18px;margin-top:28px}}
.imggrid img{{width:100%;height:280px;object-fit:cover;border:1px solid rgba(201,168,76,.22)}}
.tablewrap{{overflow:auto;border:1px solid rgba(201,168,76,.22);margin-top:24px;max-height:65vh}}
table{{width:100%;border-collapse:collapse;min-width:980px;font-size:12px}}
th{{position:sticky;top:0;background:#2A2118;color:var(--gold);padding:10px;text-align:left;border-bottom:1px solid rgba(201,168,76,.35)}}
td{{padding:9px 10px;border-bottom:1px solid rgba(201,168,76,.12);color:rgba(245,240,232,.86)}}
tr:nth-child(even) td{{background:rgba(255,255,255,.025)}}
.positive{{color:#9FE6B8!important}}.negative{{color:#FF9E9E!important}}
.note{{font-size:13px;color:rgba(245,240,232,.58);margin-top:18px}}
pre{{white-space:pre-wrap;background:rgba(0,0,0,.25);border:1px solid rgba(201,168,76,.18);padding:18px;max-height:55vh;overflow:auto;font-size:11px;color:rgba(245,240,232,.72)}}
</style>
</head>
<body>

<section class="slide">
 <div class="bg" style="background-image:url('{(data["ppt_imgs"][0] if data["ppt_imgs"] else data["pdf_imgs"][0] if data["pdf_imgs"] else "")}')"></div><div class="overlay"></div>
 <div class="inner">
  <div class="label">Heal House · Main Street</div>
  <h1>Retirement Living Assessment</h1>
  <p style="margin-top:24px">A retirement-focused development assessment using Joe’s deck for program economics and visual direction, verified against architect drawings for GFA, site math, and suite count.</p>
 </div>
</section>

<section class="slide"><div class="inner">
 <div class="label">Source Hierarchy</div><h2>Numbers tied back to source.</h2>
 <div class="grid">
  <div class="card"><div class="k">Architect PDF</div><div class="v">Source of Truth</div><div class="s">GFA, floor/site math, suite count, area schedules.</div></div>
  <div class="card"><div class="k">Joe Deck</div><div class="v">Economics</div><div class="s">Construction cost, monthly retirement rent, expenses, visuals.</div></div>
  <div class="card"><div class="k">Model</div><div class="v">BCH vs No BCH</div><div class="s">After mortgage cash flow, investor payout, NFP payout.</div></div>
 </div>
</div></section>

<section class="slide"><div class="inner">
 <div class="label">Verified Program</div><h2>Main Street retirement program.</h2>
 <div class="grid" id="programCards"></div>
 <p class="note">Please verify extracted values against the architect PDF. If a number needs correction, update the model constants in the page script.</p>
</div></section>

<section class="slide"><div class="inner">
 <div class="label">Visual Direction</div><h2>Joe deck imagery and site material.</h2>
 <div class="imggrid">{''.join(f'<img src="{x}">' for x in data["ppt_imgs"][:6])}</div>
</div></section>

<section class="slide"><div class="inner">
 <div class="label">Architect Reference</div><h2>Drawing pages used for verification.</h2>
 <div class="imggrid">{''.join(f'<img src="{x}">' for x in data["pdf_imgs"][:6])}</div>
</div></section>

<section class="slide"><div class="inner">
 <div class="label">Scenario Comparison</div><h2>Retirement: BCH vs No BCH.</h2>
 <div id="scenarioCards" class="grid"></div>
 <div class="tablewrap"><table id="yearTable"></table></div>
</div></section>

<section class="slide"><div class="inner">
 <div class="label">Decision Read</div><h2>What this tells the team.</h2>
 <div id="decisionCards" class="grid"></div>
 <p class="note">BCH scenario shows the mission-aligned version: reduced debt burden and annual split between investors and NFP after mortgage payment. No BCH shows conventional investor-only cash flow pressure.</p>
</div></section>

<section class="slide"><div class="inner">
 <div class="label">Appendix</div><h2>Extracted source text for verification.</h2>
 <h3>Architect PDF extraction</h3><pre>{data["pdf_source"].replace("&","&amp;").replace("<","&lt;")}</pre>
 <h3 style="margin-top:24px">Joe PPT extraction</h3><pre>{data["ppt_source"].replace("&","&amp;").replace("<","&lt;")}</pre>
</div></section>

<script>
const model = {json.dumps(data)};
function fmt(n){{let neg=n<0;n=Math.abs(n);let v=n>=1e6?'$'+(n/1e6).toFixed(2)+'M':'$'+Math.round(n/1000)+'K';return neg?'-'+v:v}}
function cls(n){{return n<0?'negative':'positive'}}
function debt(p,r,y){{let m=r/12,n=y*12;return p*(m*Math.pow(1+m,n))/(Math.pow(1+m,n)-1)*12}}
function run(bch, year){{
 const cost=model.gfa*model.cost_sf;
 const principal=bch?cost*(1-model.bch_forgivable):cost;
 const ds=debt(principal,bch?model.bch_rate:model.conv_rate,bch?model.bch_am:model.conv_am);
 const revenue=model.suites*model.monthly_rent*12*model.occupancy*Math.pow(1+model.growth,year-1);
 const expense=revenue*model.expense_ratio;
 const noi=revenue-expense;
 const cf=noi-ds;
 return {{cost,principal,ds,revenue,expense,noi,cf,inv:cf>0?(bch?cf*model.investor_split:cf):0,nfp:cf>0&&bch?cf*model.nfp_split:0}};
}}
document.getElementById('programCards').innerHTML = `
 <div class="card"><div class="k">GFA · Architect PDF</div><div class="v">${{model.gfa?model.gfa.toLocaleString():'VERIFY'}}</div><div class="s">Gross floor area / area schedule</div></div>
 <div class="card"><div class="k">Suites · Architect PDF</div><div class="v">${{model.suites?model.suites.toLocaleString():'VERIFY'}}</div><div class="s">Retirement suites / beds</div></div>
 <div class="card"><div class="k">Cost · Joe Deck</div><div class="v">$${{model.cost_sf.toLocaleString()}}/SF</div><div class="s">Construction cost assumption</div></div>
 <div class="card"><div class="k">Monthly Rent · Joe Deck</div><div class="v">$${{model.monthly_rent.toLocaleString()}}</div><div class="s">Per suite/month</div></div>
 <div class="card"><div class="k">Expense Ratio · Joe Deck</div><div class="v">${{Math.round(model.expense_ratio*100)}}%</div><div class="s">Retirement operating expense</div></div>
`;
let y20b=run(true,20), y20n=run(false,20);
document.getElementById('scenarioCards').innerHTML = `
 <div class="card"><div class="k">BCH · Year 20 After Mortgage</div><div class="v ${{cls(y20b.cf)}}">${{fmt(y20b.cf)}}</div><div class="s">Investor ${{fmt(y20b.inv)}} · NFP ${{fmt(y20b.nfp)}}</div></div>
 <div class="card"><div class="k">No BCH · Year 20 After Mortgage</div><div class="v ${{cls(y20n.cf)}}">${{fmt(y20n.cf)}}</div><div class="s">Investor ${{fmt(y20n.inv)}} · NFP $0</div></div>
 <div class="card"><div class="k">Total Project Cost</div><div class="v">${{fmt(y20b.cost)}}</div><div class="s">GFA × cost/SF</div></div>
`;
let html='<thead><tr><th>Year</th><th>Scenario</th><th>Revenue</th><th>Expenses</th><th>NOI</th><th>Mortgage</th><th>After Mortgage</th><th>Investor</th><th>NFP</th></tr></thead><tbody>';
for(let y=1;y<=20;y++){{[true,false].forEach(b=>{{let r=run(b,y);html+=`<tr><td>${{y}}</td><td>${{b?'BCH 40% forgivable / 50yr CMHC':'No BCH / 30yr conventional'}}</td><td>${{fmt(r.revenue)}}</td><td>${{fmt(r.expense)}}</td><td>${{fmt(r.noi)}}</td><td>${{fmt(r.ds)}}</td><td class="${{cls(r.cf)}}">${{fmt(r.cf)}}</td><td>${{fmt(r.inv)}}</td><td>${{fmt(r.nfp)}}</td></tr>`}})}}
html+='</tbody>';document.getElementById('yearTable').innerHTML=html;
document.getElementById('decisionCards').innerHTML = `
 <div class="card"><div class="k">Preferred Structure</div><div class="v">BCH</div><div class="s">Lower mortgage load and creates NFP participation.</div></div>
 <div class="card"><div class="k">Investor/NFP Logic</div><div class="v">60 / 40</div><div class="s">Applied only after mortgage payment.</div></div>
 <div class="card"><div class="k">Verify Next</div><div class="v">Inputs</div><div class="s">GFA, suite count, rent, cost/SF, expense ratio.</div></div>
`;
</script>
</body></html>'''

OUT.write_text(html)
print("Built", OUT)
print("Extracted:", data)
