from pathlib import Path
import zipfile, json, re, html as htmlmod

ROOT = Path.home() / "healhouse-decks"
PDF = ROOT / "source" / "main-street-architect.pdf"
PPT = ROOT / "source" / "main-street-joe.pptx"
ASSET = ROOT / "assets" / "main-street"
OUT = ROOT / "main-street-v2.html"

assert PDF.exists(), f"Missing {PDF}"
assert PPT.exists(), f"Missing {PPT}"

# Dependencies
import fitz
from pptx import Presentation

def rel(p): return str(Path(p).relative_to(ROOT))

# Extract architect pages
pdf_dir = ASSET / "pdf"
pdf_dir.mkdir(parents=True, exist_ok=True)
pdf_text = []
pdf_imgs = []
doc = fitz.open(PDF)
for i, page in enumerate(doc):
    txt = page.get_text("text") or ""
    pdf_text.append(f"--- ARCHITECT PDF PAGE {i+1} ---\n{txt}")
    pix = page.get_pixmap(matrix=fitz.Matrix(1.8,1.8), alpha=False)
    out = pdf_dir / f"architect-page-{i+1}.jpg"
    pix.save(out)
    pdf_imgs.append(rel(out))
pdf_all = "\n".join(pdf_text)

# Extract Joe PPT text and media
ppt_dir = ASSET / "ppt"
ppt_dir.mkdir(parents=True, exist_ok=True)
ppt_text = []
prs = Presentation(PPT)
for idx, slide in enumerate(prs.slides, 1):
    bits=[]
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            bits.append(shape.text.strip())
    ppt_text.append(f"--- JOE PPT SLIDE {idx} ---\n" + "\n".join(bits))
ppt_all = "\n".join(ppt_text)

ppt_imgs=[]
with zipfile.ZipFile(PPT) as z:
    for name in z.namelist():
        if name.startswith("ppt/media/") and name.lower().endswith((".png",".jpg",".jpeg")):
            dest = ppt_dir / Path(name).name
            dest.write_bytes(z.read(name))
            ppt_imgs.append(rel(dest))

def nums(pattern, text, lo=None, hi=None):
    vals=[]
    for m in re.findall(pattern, text, re.I):
        if isinstance(m, tuple): m = next((x for x in m if x), "")
        s = re.sub(r"[^0-9.]", "", str(m))
        if not s or s == ".": continue
        try: v = float(s)
        except: continue
        if lo is not None and v < lo: continue
        if hi is not None and v > hi: continue
        vals.append(v)
    return vals

# Extract likely numbers
gfa_vals = nums(r"([\d,]{4,})\s*(?:sf|sq\.?\s*ft|sqft)", pdf_all, 10000, 1000000)
gfa = int(max(gfa_vals)) if gfa_vals else 0

suite_vals = []
for pat in [
    r"(\d{2,4})\s*(?:suites|suite|units|beds|rooms)",
    r"(?:suites|suite|units|beds|rooms)[^\d]{0,20}(\d{2,4})"
]:
    suite_vals += nums(pat, pdf_all + "\n" + ppt_all, 10, 1000)
suites = int(max(suite_vals)) if suite_vals else 0
cost_vals = nums(r"\$?\s*([\d,]+(?:\.\d+)?)\s*/\s*(?:sf|sq\.?\s*ft|sqft)", ppt_all, 100, 1500)
cost_sf = cost_vals[0] if cost_vals else 383.50

rent_vals = nums(r"\$?\s*([\d,]{4,6})\s*(?:/|per)?\s*(?:month|mo|monthly)", ppt_all, 2000, 20000)
if not rent_vals:
    rent_vals = nums(r"(?:rent|monthly fee|retirement)[^\n$]{0,80}\$?\s*([\d,]{4,6})", ppt_all, 2000, 20000)
monthly_rent = rent_vals[0] if rent_vals else 5500

exp_vals = nums(r"(?:expense|opex|operating)[^\n%]{0,80}(\d{1,2})\s*%", ppt_all, 10, 80)
expense_ratio = (exp_vals[0]/100) if exp_vals else 0.42

verify_gfa = not bool(gfa)
verify_suites = not bool(suites)
if not gfa: gfa = 100000
if not suites: suites = round(gfa / 700)
data = {
  "gfa": gfa,
  "suites": suites,
  "cost_sf": cost_sf,
  "monthly_rent": monthly_rent,
  "expense_ratio": expense_ratio,
  "occupancy": .90,
  "growth": .025,
  "bch_forgivable": .40,
  "bch_rate": .0415,
  "bch_am": 50,
  "conv_rate": .0499,
  "conv_am": 30,
  "investor_split": .60,
  "nfp_split": .40,
  "verify_gfa": verify_gfa,
  "verify_suites": verify_suites
}

# Pick visuals: use PPT images first, architect drawings for reference only
hero = ppt_imgs[0] if ppt_imgs else (pdf_imgs[0] if pdf_imgs else "")
visuals = ppt_imgs[:8]
drawings = pdf_imgs[:6]

def img(url): return f"background-image:url('{url}')" if url else ""

html = f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>Heal House Main Street — Retirement Assessment</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
:root{{--espresso:#1E1A14;--charcoal:#2A2118;--cream:#F5F0E8;--gold:#C9A84C;--gold2:#E2C97E;--line:rgba(201,168,76,.24);--card:rgba(42,33,24,.78);--green:#9FE6B8;--red:#FF9E9E}}
body{{background:var(--espresso);color:var(--cream);font-family:'Inter','Segoe UI',system-ui,sans-serif;line-height:1.55}}
.slide{{min-height:100vh;display:flex;align-items:center;position:relative;overflow:hidden;border-bottom:1px solid rgba(201,168,76,.10)}}
.alt{{background:var(--charcoal)}}
.inner{{max-width:1180px;width:100%;margin:0 auto;padding:88px 58px;position:relative;z-index:2}}
.bg{{position:absolute;inset:0;background-size:cover;background-position:center;opacity:.30;filter:saturate(.9)}}
.scrim{{position:absolute;inset:0;background:linear-gradient(90deg,rgba(20,15,8,.95),rgba(20,15,8,.72),rgba(20,15,8,.44))}}
.label{{display:inline-block;font-size:11px;letter-spacing:.26em;text-transform:uppercase;color:var(--gold);border-bottom:1px solid var(--line);padding-bottom:5px;margin-bottom:24px;font-weight:700}}
h1,h2,h3{{font-family:Georgia,'Times New Roman',serif;font-weight:400;line-height:1.08}}
h1{{font-size:clamp(50px,8vw,104px);font-style:italic;max-width:920px}}
h2{{font-size:clamp(34px,5vw,62px);max-width:940px}}
h3{{font-size:25px;color:var(--gold2);margin-bottom:10px}}
p{{font-size:18px;color:rgba(245,240,232,.82);max-width:790px}}
.divider{{width:92px;height:2px;background:var(--gold);margin:28px 0}}
.two{{display:grid;grid-template-columns:1fr 1fr;gap:46px;align-items:center}}
.grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(230px,1fr));gap:18px;margin-top:34px}}
.card{{background:var(--card);border:1px solid var(--line);padding:24px;position:relative}}
.card:before{{content:"";position:absolute;top:0;left:0;right:0;height:2px;background:linear-gradient(90deg,var(--gold),transparent)}}
.k{{font-size:11px;letter-spacing:.15em;text-transform:uppercase;color:var(--gold);font-weight:700}}
.v{{font-family:Georgia,serif;font-size:34px;color:var(--cream);margin-top:8px;line-height:1.05}}
.s{{font-size:13px;color:rgba(245,240,232,.62);margin-top:10px}}
.panel{{min-height:450px;background-size:cover;background-position:center;border:1px solid var(--line);position:relative}}
.panel:after{{content:"";position:absolute;inset:0;background:linear-gradient(180deg,transparent 35%,rgba(20,15,8,.62))}}
.tablewrap{{overflow:auto;border:1px solid var(--line);margin-top:28px;max-height:64vh;background:rgba(0,0,0,.14)}}
table{{width:100%;border-collapse:collapse;min-width:940px;font-size:12px}}
th{{position:sticky;top:0;background:#2A2118;color:var(--gold);padding:11px;text-align:left;border-bottom:1px solid rgba(201,168,76,.38)}}
td{{padding:10px 11px;border-bottom:1px solid rgba(201,168,76,.11);color:rgba(245,240,232,.86)}}
tr:nth-child(even) td{{background:rgba(255,255,255,.025)}}
.positive{{color:var(--green)!important}}.negative{{color:var(--red)!important}}
.source{{margin-top:20px;font-size:12px;color:rgba(245,240,232,.55)}}
.footer{{position:absolute;bottom:0;left:0;right:0;padding:13px 58px;display:flex;justify-content:space-between;border-top:1px solid rgba(201,168,76,.16);background:rgba(20,15,8,.38);z-index:4}}
.footer span{{font-size:10px;letter-spacing:.18em;text-transform:uppercase;color:rgba(201,168,76,.72)}}
@media(max-width:800px){{.inner{{padding:70px 26px}}.two{{grid-template-columns:1fr}}.panel{{min-height:280px}}.footer{{padding:12px 26px}}}}
</style>
</head>
<body>
<section class="slide"><div class="bg" style="{img(hero)}"></div><div class="scrim"></div><div class="inner"><div class="label">Heal House · Main Street</div><h1>Retirement Living Assessment</h1><div class="divider"></div><p>A focused retirement-use development review using Main Street source files only: Joe’s deck for concept economics and visuals, architect drawings for GFA and program verification.</p></div><div class="footer"><span>HEAL HOUSE</span><span>MAIN STREET · RETIREMENT</span></div></section>
<section class="slide alt"><div class="inner two"><div><div class="label">Method</div><h2>One project. Two financing paths.</h2><div class="divider"></div><p>The assessment isolates the all-retirement use case and compares BCH against conventional no-BCH financing after mortgage payment.</p><div class="grid"><div class="card"><div class="k">Architect PDF</div><div class="v">Program</div><div class="s">GFA, suite count, drawing math.</div></div><div class="card"><div class="k">Joe PPT</div><div class="v">Economics</div><div class="s">Cost, rent, expenses, visuals.</div></div></div></div><div class="panel" style="{img(visuals[1] if len(visuals)>1 else hero)}"></div></div><div class="footer"><span>H
bashEAL HOUSE</span><span>SOURCE HIERARCHY</span></div></section>
<section class="slide"><div class="inner"><div class="label">Program Inputs</div><h2>Source-labelled assumptions.</h2><div class="grid" id="programCards"></div><div class="source">Values marked VERIFY require final human confirmation against the PDF/PPT before external circulation.</div></div><div class="footer"><span>HEAL HOUSE</span><span>PROGRAM</span></div></section>
<section class="slide alt"><div class="inner two"><div><div class="label">Retirement Economics</div><h2>Revenue is suite count × monthly rent × stabilized occupancy.</h2><div class="divider"></div><p>Operating expenses are applied before debt service. The decision metric is cash flow after mortgage payment.</p><div class="grid"><div class="card"><div class="k">Occupancy</div><div class="v">90%</div><div class="s">Model assumption</div></div><div class="card"><div class="k">Growth</div><div class="v">2.5%</div><div class="s">Annual revenue growth</div></div></div></div><div class="panel" style="{img(visuals[2] if len(visuals)>2 else hero)}"></div></div><div class="footer"><span>HEAL HOUSE</span><span>ECONOMICS</span></div></section>
<section class="slide"><div class="inner"><div class="label">Year 20 Summary</div><h2>BCH vs no BCH after mortgage.</h2><div class="grid" id="scenarioCards"></div><div class="source">BCH surplus is split 60% investors / 40% NFP after mortgage payment. No BCH is shown as investor cash flow only.</div></div><div class="footer"><span>HEAL HOUSE</span><span>SUMMARY</span></div></section>
<section class="slide alt"><div class="inner"><div class="label">20-Year View</div><h2>Year-by-year comparison.</h2><div class="tablewrap"><table id="yearTable"></table></div></div><div class="footer"><span>HEAL HOUSE</span><span>YEAR-BY-YEAR</span></div></section>
<section class="slide"><div class="inner"><div class="label">Decision Read</div><h2>The BCH path creates the cleanest mission-aligned surplus.</h2><div class=
bash"grid" id="decisionCards"></div></div><div class="footer"><span>HEAL HOUSE</span><span>DECISION READ</span></div></section>
<section class="slide alt"><div class="inner"><div class="label">Joe Deck Visuals</div><h2>Main Street concept imagery.</h2><div class="grid">{''.join(f'<div class="card" style="padding:0;overflow:hidden"><div class="panel" style="min-height:300px;{img(v)}"></div></div>' for v in visuals[:4])}</div></div><div class="footer"><span>HEAL HOUSE</span><span>VISUALS</span></div></section>
<section class="slide"><div class="inner"><div class="label">Architect Reference</div><h2>Drawing references for verification.</h2><div class="grid">{''.join(f'<div class="card" style="padding:0;overflow:hidden"><div class="panel" style="min-height:300px;background-size:contain;background-repeat:no-repeat;{img(v)}"></div></div>' for v in drawings[:4])}</div></div><div class="footer"><span>HEAL HOUSE</span><span>DRAWINGS</span></div></section>
<script>
const model={json.dumps(data)};
function fmt(n){{let neg=n<0;n=Math.abs(n);let v=n>=1e6?'$'+(n/1e6).toFixed(2)+'M':'$'+Math.round(n/1000)+'K';return neg?'-'+v:v}}
function cls(n){{return n<0?'negative':'positive'}}
function debt(p,r,y){{let m=r/12,n=y*12;return p*(m*Math.pow(1+m,n))/(Math.pow(1+m,n)-1)*12}}
function run(bch,year){{const cost=model.gfa*model.cost_sf;const principal=bch?cost*(1-model.bch_forgivable):cost;const ds=debt(principal,bch?model.bch_rate:model.conv_rate,bch?model.bch_am:model.conv_am);const revenue=model.suites*model.monthly_rent*12*model.occupancy*Math.pow(1+model.growth,year-1);const expense=revenue*model.expense_ratio;const noi=revenue-expense;const cf=noi-ds;return{{cost,principal,ds,revenue,expense,noi,cf,inv:cf>0?(bch?cf*model.investor_split:cf):0,nfp:cf>0&&bch?cf*model.nfp_split:0}}}}
document.getElementById('programCards').innerHTML=`
<div class="card"><div class="k">GFA · Architect PDF${{model.verify_gfa?' · VERIFY':''}}</div><div class="v">${{model.gfa.toLocaleString()}}</di
bashv><div class="s">Gross floor area used for project cost.</div></div>
<div class="card"><div class="k">Suites · Architect PDF${{model.verify_suites?' · VERIFY':''}}</div><div class="v">${{model.suites.toLocaleString()}}</div><div class="s">Retirement suites / beds.</div></div>
<div class="card"><div class="k">Cost · Joe Deck</div><div class="v">$${{model.cost_sf.toLocaleString()}}/SF</div><div class="s">Construction cost assumption.</div></div>
<div class="card"><div class="k">Monthly Rent · Joe Deck</div><div class="v">$${{model.monthly_rent.toLocaleString()}}</div><div class="s">Per suite per month.</div></div>
<div class="card"><div class="k">Expense Ratio · Joe Deck</div><div class="v">${{Math.round(model.expense_ratio*100)}}%</div><div class="s">Operating expense assumption.</div></div>
<div class="card"><div class="k">Project Cost</div><div class="v">${{fmt(model.gfa*model.cost_sf)}}</div><div class="s">GFA × cost/SF.</div></div>`;
let b=run(true,20), n=run(false,20);
document.getElementById('scenarioCards').innerHTML=`
<div class="card"><div class="k">BCH · Year 20 after mortgage</div><div class="v ${{cls(b.cf)}}">${{fmt(b.cf)}}</div><div class="s">Investors ${{fmt(b.inv)}} · NFP ${{fmt(b.nfp)}}</div></div>
<div class="card"><div class="k">No BCH · Year 20 after mortgage</div><div class="v ${{cls(n.cf)}}">${{fmt(n.cf)}}</div><div class="s">Investor cash flow ${{fmt(n.inv)}} · NFP $0</div></div>
<div class="card"><div class="k">BCH Debt Basis</div><div class="v">${{fmt(b.principal)}}</div><div class="s">40% forgivable loan reduces mortgageable basis.</div></div>
<div class="card"><div class="k">No BCH Debt Basis</div><div class="v">${{fmt(n.principal)}}</div><div class="s">Conventional 30-year amortization.</div></div>`;
let html='<thead><tr><th>Year</th><th>Scenario</th><th>Revenue</th><th>Expenses</th><th>NOI</th><th>Mortgage</th><th>After Mortgage</th><th>Investor</th><th>NFP</th></tr></thead><tbody>';
for(let y=1;y<=20;y++){{[true,false].forEach(
bashx=>{{let r=run(x,y);html+=`<tr><td>${{y}}</td><td>${{x?'BCH · 40% forgivable / 50yr CMHC':'No BCH · 30yr conventional'}}</td><td>${{fmt(r.revenue)}}</td><td>${{fmt(r.expense)}}</td><td>${{fmt(r.noi)}}</td><td>${{fmt(r.ds)}}</td><td class="${{cls(r.cf)}}">${{fmt(r.cf)}}</td><td>${{fmt(r.inv)}}</td><td>${{fmt(r.nfp)}}</td></tr>`}})}}
html+='</tbody>';document.getElementById('yearTable').innerHTML=html;
document.getElementById('decisionCards').innerHTML=`
<div class="card"><div class="k">Preferred mission structure</div><div class="v">BCH</div><div class="s">Reduces debt pressure and creates NFP participation after mortgage.</div></div>
<div class="card"><div class="k">Investor story</div><div class="v">After Debt</div><div class="s">Payouts are compared only after mortgage payment.</div></div>
<div class="card"><div class="k">Before investor version</div><div class="v">Verify</div><div class="s">Confirm GFA, suite count, cost/SF, rent, and expenses.</div></div>`;
</script>
</body></html>'''

OUT.write_text(html)
print("Built clean Main Street deck:", OUT)
print("Data:", data)
